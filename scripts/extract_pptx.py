"""Extract slide text, notes, tables, and images from PPTX into markdown.

Usage:
	python extract_pptx.py <input.pptx> <output_dir>
"""

from __future__ import annotations

import io
import re
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation


def usage() -> None:
	print("Usage: python extract_pptx.py <input.pptx> <output_dir>")


def sanitize(value: str) -> str:
	cleaned = re.sub(r"[^A-Za-z0-9._ -]", "_", value.strip())
	return cleaned or "presentation"


def escape_md(value: str) -> str:
	return value.replace("|", "\\|").replace("\n", " ").strip()


def rows_to_markdown(rows: list[list[str]]) -> str:
	if not rows:
		return ""
	width = max(len(row) for row in rows)
	padded = [row + [""] * (width - len(row)) for row in rows]
	header = [escape_md(cell) or " " for cell in padded[0]]
	divider = ["---"] * width

	lines = [
		"| " + " | ".join(header) + " |",
		"| " + " | ".join(divider) + " |",
	]
	for row in padded[1:]:
		lines.append("| " + " | ".join(escape_md(cell) or " " for cell in row) + " |")
	return "\n".join(lines)


def save_blob(blob: bytes, out_path: Path) -> None:
	try:
		with Image.open(io.BytesIO(blob)) as image:
			image.save(out_path)
	except Exception:
		out_path.write_bytes(blob)


def main() -> int:
	if len(sys.argv) != 3:
		usage()
		return 1

	input_path = Path(sys.argv[1]).expanduser().resolve()
	output_dir = Path(sys.argv[2]).expanduser().resolve()

	if not input_path.exists() or input_path.suffix.lower() != ".pptx":
		print(f"Invalid PPTX input file: {input_path}")
		return 1

	output_dir.mkdir(parents=True, exist_ok=True)
	images_dir = output_dir / "images"
	images_dir.mkdir(parents=True, exist_ok=True)
	markdown_path = output_dir / f"{sanitize(input_path.stem)}.md"

	presentation = Presentation(input_path)
	lines: list[str] = [f"# {input_path.name}", "", f"Slides: {len(presentation.slides)}", ""]

	for slide_index, slide in enumerate(presentation.slides, start=1):
		lines.append(f"## Slide {slide_index}")
		lines.append("")

		slide_text: list[str] = []
		tables_found = 0
		image_count = 0

		for shape in slide.shapes:
			if hasattr(shape, "text") and shape.text:
				text = shape.text.strip()
				if text:
					slide_text.append(text)

			if shape.has_table:
				tables_found += 1
				rows = []
				for row in shape.table.rows:
					rows.append([cell.text.strip() for cell in row.cells])
				lines.append(f"### Table {slide_index}.{tables_found}")
				lines.append("")
				lines.append(rows_to_markdown(rows))
				lines.append("")

			if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
				image = shape.image
				ext = image.ext or "png"
				image_count += 1
				image_name = f"slide_{slide_index:03d}_image_{image_count:02d}.{ext}"
				image_path = images_dir / image_name
				save_blob(image.blob, image_path)
				lines.append(f"- ![Slide {slide_index} image {image_count}](./images/{image_name})")

		if slide_text:
			lines.append("### Text")
			lines.append("")
			for paragraph in slide_text:
				lines.append(f"- {paragraph}")
			lines.append("")
		else:
			lines.append("_No text extracted from slide body._")
			lines.append("")

		notes_text = ""
		if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
			notes_text = slide.notes_slide.notes_text_frame.text.strip()

		lines.append("### Speaker Notes")
		lines.append("")
		lines.append(notes_text if notes_text else "_No speaker notes found._")
		lines.append("")

	markdown_path.write_text("\n".join(lines), encoding="utf-8")
	print(f"Extracted: {input_path}")
	print(f"Markdown: {markdown_path}")
	print(f"Images: {images_dir}")
	return 0


if __name__ == "__main__":
	raise SystemExit(main())
