"""
Extract text, tables, speaker notes, and images from PowerPoint (.pptx) files.
Renders full slides as images (capturing diagrams, SmartArt, charts, grouped shapes)
via PowerPoint COM automation on Windows, with fallback to embedded-image-only extraction.

Usage: python extract_pptx.py <input.pptx> <output_dir>
"""

import sys
import os
from pathlib import Path

# Minimum image file size to keep (bytes). Filters out tiny icons/decorative elements.
MIN_IMAGE_SIZE = 10_000


def render_slides_via_com(input_path: str, images_dir: Path, slide_count: int) -> bool:
    """Render every slide as a high-res PNG using PowerPoint COM automation.
    Returns True if successful, False if PowerPoint is unavailable."""
    abs_path = str(Path(input_path).resolve())
    try:
        import comtypes.client

        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        # Open read-only, no window
        presentation = powerpoint.Presentations.Open(
            abs_path, ReadOnly=True, Untitled=False, WithWindow=False
        )

        for i in range(1, slide_count + 1):
            out_file = str(images_dir / f"slide{i}_full.png")
            presentation.Slides(i).Export(out_file, "PNG", 1920, 1080)

        presentation.Close()
        powerpoint.Quit()
        print(f"Rendered {slide_count} slides via PowerPoint COM")
        return True

    except Exception as e:
        print(f"PowerPoint COM rendering unavailable ({e}), trying win32com...")

    try:
        import win32com.client

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            abs_path, ReadOnly=True, Untitled=False, WithWindow=False
        )

        for i in range(1, slide_count + 1):
            out_file = str(images_dir / f"slide{i}_full.png")
            presentation.Slides(i).Export(out_file, "PNG", 1920, 1080)

        presentation.Close()
        powerpoint.Quit()
        print(f"Rendered {slide_count} slides via win32com")
        return True

    except Exception as e:
        print(f"win32com rendering also unavailable ({e})")
        return False


def extract_pptx(input_path: str, output_dir: str) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(input_path)
    output_path = Path(output_dir)
    images_dir = output_path / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    filename = Path(input_path).stem
    slide_count = len(prs.slides)

    # --- Phase 1: Render full slides as images (captures diagrams, SmartArt, etc.) ---
    slides_rendered = render_slides_via_com(input_path, images_dir, slide_count)

    # --- Phase 2: Extract text, tables, embedded images, and speaker notes ---
    md_lines = [f"# Extracted: {filename}\n"]
    md_lines.append(f"**Source**: {Path(input_path).name}\n")
    md_lines.append(f"**Slides**: {slide_count}\n")
    if slides_rendered:
        md_lines.append("**Slide rendering**: Full slides rendered via PowerPoint\n")
    else:
        md_lines.append("**Slide rendering**: Embedded images only (PowerPoint not available for full rendering)\n")
    md_lines.append("---\n")

    embedded_count = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        md_lines.append(f"\n## Slide {slide_num}\n")

        # Embed the full-slide render if available
        full_slide_img = f"slide{slide_num}_full.png"
        if slides_rendered and (images_dir / full_slide_img).exists():
            md_lines.append(f"![Slide {slide_num}](./images/{full_slide_img})\n")

        # Extract title
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                md_lines.append(f"### {title_text}\n")

        # Extract text from all shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text and (not slide.shapes.title or text != slide.shapes.title.text.strip()):
                        level = paragraph.level if paragraph.level else 0
                        indent = "  " * level
                        md_lines.append(f"{indent}- {text}")

            # Extract tables
            if shape.has_table:
                table = shape.table
                md_lines.append("")
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                md_lines.append("| " + " | ".join(headers) + " |")
                md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row_idx, row in enumerate(table.rows):
                    if row_idx == 0:
                        continue
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    md_lines.append("| " + " | ".join(cells) + " |")
                md_lines.append("")

            # Extract embedded high-res images (skip tiny icons)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                blob = image.blob
                if len(blob) >= MIN_IMAGE_SIZE:
                    embedded_count += 1
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    image_name = f"slide{slide_num}_embed{embedded_count}.{ext}"
                    image_path = images_dir / image_name
                    with open(image_path, "wb") as f:
                        f.write(blob)
                    alt_text = shape.name or f"Embedded image from slide {slide_num}"
                    md_lines.append(f"\n![{alt_text}](./images/{image_name})\n")

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                md_lines.append(f"\n> **Speaker Notes**: {notes_text}\n")

        md_lines.append("\n---\n")

    # Write markdown
    md_path = output_path / f"{filename}.md"
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))

    rendered_msg = f", {slide_count} full-slide renders" if slides_rendered else " (no full-slide renders)"
    print(f"Extracted {slide_count} slides, {embedded_count} embedded images{rendered_msg}")
    print(f"Output: {md_path}")
    return str(md_path)


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python extract_pptx.py <input.pptx> <output_dir>")
        sys.exit(1)

    input_file = sys.argv[1]
    output_directory = sys.argv[2]

    if not os.path.exists(input_file):
        print(f"Error: File not found: {input_file}")
        sys.exit(1)

    extract_pptx(input_file, output_directory)
